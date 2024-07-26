from pptx import Presentation
import os
from espnet_model_zoo.downloader import ModelDownloader
from espnet2.bin.tts_inference import Text2Speech
from espnet2.utils.types import str_or_none
import soundfile
import torch
import logging
logger = logging.getLogger('pptx2video.pptx2audio')

def extract_note_description(tagetfile):
    presentation = Presentation(tagetfile)
    # note extraction
    notes ={}
    i = 0
    for slide in presentation.slides:
        note_text =""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            logger.debug(slide.notes_slide.notes_text_frame.text)
            note_text = slide.notes_slide.notes_text_frame.text
            logger.debug('-----')
        notes[i] = note_text
        i += 1
    return notes

def make_wav(text2speech,filepath,text):

    if text==None or text=="":
        return
    try:
        logger.debug("Speech generation start " + filepath )
        with torch.no_grad():
            wav = text2speech(text)["wav"]
        wavdata = wav.view(-1).cpu().numpy()
        samplerate=text2speech.fs
        soundfile.write(filepath, wavdata, samplerate, subtype='PCM_24')    
        logger.debug("Speech generation end " + filepath )
    except Exception as e:
        print(e)
        return

def pptx_note2_audio(tag,tagetfile,audiodir,device='cpu'):

    lang = 'Japanese'
    model_path  ="/root/src/models/tts_finetune_jvs010_jsut_vits_raw_phn_jaconv_pyopenjtalk_prosody_latest/exp/tts_finetune_jvs010_jsut_vits_raw_phn_jaconv_pyopenjtalk_prosody/100epoch.pth"
    train_config  ="/root/src/models/tts_finetune_jvs010_jsut_vits_raw_phn_jaconv_pyopenjtalk_prosody_latest/exp/tts_finetune_jvs010_jsut_vits_raw_phn_jaconv_pyopenjtalk_prosody/config.yaml"
    vocoder_tag = 'none'
    text2speech = None
    if tag:
        text2speech = Text2Speech.from_pretrained(
            model_tag=str_or_none(tag),
            vocoder_tag=str_or_none(vocoder_tag),
            device=device
        )
    else: 
        text2speech = Text2Speech.from_pretrained(
            model_file=model_path, 
            train_config=train_config,
            vocoder_tag=str_or_none(vocoder_tag),
            device=device
        )

    notes = extract_note_description(tagetfile)
    output_audio_files = {}
    for i in range(len(notes)):
        audio_path = audiodir+"/audio_"+str(i) +".wav"
        make_wav(text2speech,audio_path,str(notes[i]))
        output_audio_files[i] = audio_path

    return output_audio_files,notes

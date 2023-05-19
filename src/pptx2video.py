import sys

# power point
from pptx2jpg import *
from pptx import Presentation
import glob
import re
import os
import tempfile

# Video 
import cv2
import numpy as np
import ffmpeg
import torch

# audio
import soundfile
from espnet_model_zoo.downloader import ModelDownloader
from espnet2.bin.tts_inference import Text2Speech
from espnet2.utils.types import str_or_none
from pydub import AudioSegment
import math
from pydub import AudioSegment
from moviepy.editor import *

device = "cpu"
if torch.cuda.is_available():
    device ="cuda"
    print(torch.cuda.get_device_name())

lang = 'Japanese'
model_path  ="/root/src/models/tts_finetune_jvs010_jsut_vits_raw_phn_jaconv_pyopenjtalk_prosody_latest/exp/tts_finetune_jvs010_jsut_vits_raw_phn_jaconv_pyopenjtalk_prosody/100epoch.pth"
train_config  ="/root/src/models/tts_finetune_jvs010_jsut_vits_raw_phn_jaconv_pyopenjtalk_prosody_latest/exp/tts_finetune_jvs010_jsut_vits_raw_phn_jaconv_pyopenjtalk_prosody/config.yaml"
vocoder_tag = 'none'
text2speech = Text2Speech.from_pretrained(
    model_file=model_path, 
    train_config=train_config,
    vocoder_tag=str_or_none(vocoder_tag),
    device=device
)

def make_wav(filepath,text):

    if text==None or text=="":
        return
    try:
        print("Speech generation start " + filepath )
        with torch.no_grad():
            wav = text2speech(text)["wav"]
        wavdata = wav.view(-1).cpu().numpy()
        samplerate=text2speech.fs
        soundfile.write(filepath, wavdata, samplerate, subtype='PCM_24')    
        print("Speech generation end " + filepath )
    except Exception as e:
        print(e)
        return

def pptx_note2_audio(tagetfile,audiodir):
    presentation = Presentation(tagetfile)

    # note extraction
    notes ={}
    i = 0
    for slide in presentation.slides:
        note_text =""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            print(slide.notes_slide.notes_text_frame.text)
            note_text = slide.notes_slide.notes_text_frame.text
            print('-----')
        notes[i] = note_text
        i += 1

    output_audio_files = {}
    for i in range(len(notes)):
        audio_path = audiodir+"/audio_"+str(i) +".wav"
        make_wav(audio_path,str(notes[i]))
        output_audio_files[i] = audio_path

    return output_audio_files,notes

def searchfiles(targetdir,ext):
    def atoi(text):
        return int(text) if text.isdigit() else text

    def natural_keys(text):
        return [ atoi(c) for c in re.split(r'(\d+)', text) ]

    files = sorted(glob.glob(targetdir + "/*" + ext), key=natural_keys)

    for file in files:
        print(file)

    return files

def jpg2video(tagetfile,imagedir,outputpath,slide_time):

    output_audio_files,notes = pptx_note2_audio(tagetfile,imagedir)
    pptx2jpg(tagetfile,imagedir)
    files = searchfiles(imagedir,".jpg")

    outputvideopath = imagedir + '/output.mp4'
    outputvideotmppath = imagedir +'/output_tmp.mp4'

    fourcc = cv2.VideoWriter_fourcc('m','p','4', 'v')
    w=1270
    h=720
    fps = 20
    video  = cv2.VideoWriter(outputvideotmppath, fourcc,fps, (w, h))

    audio_index =0
    final_sound = None
    for img_file in files:
        print("read:" + img_file)
        img = cv2.imread(img_file)
        img1 = cv2.resize( img, (w, h) ) 
        if img1 is None:
            print("can't read")
            break    

        frame_num = fps*slide_time
        # Adjust the number of frames to match the audio content
        audio_path = output_audio_files[audio_index]
        if audio_path and os.path.exists(audio_path):
            sound1 = AudioSegment.from_file(audio_path)
            print(str(sound1.duration_seconds))
            padding = slide_time - sound1.duration_seconds
            if padding > 0.0:
                sound1 += AudioSegment.silent(duration=padding*1000)
            else:
                frame_num += int(fps*abs(padding))
            if final_sound == None:
                final_sound = sound1
            else:
                final_sound += sound1
        else:
            if final_sound== None:
                final_sound = AudioSegment.silent(duration=slide_time*1000)
            else:
                final_sound += AudioSegment.silent(duration=slide_time*1000)

        #final_sound += AudioSegment.silent(duration=1000)
        #frame_num += fps
        print("frame count:" + str(frame_num) )
        for i in range(frame_num):
            video.write(img1)
        audio_index+=1
    video.release()
    convertfile = imagedir + "/audio_conv.wav"
    final_sound.export(convertfile, format='wav')

    print(outputvideotmppath)
    # H.264 
    ffmpeg.input(outputvideotmppath).output(outputvideopath, vcodec='libx264').run(overwrite_output=True)
    os.remove(outputvideotmppath)

    video_clip = VideoFileClip(outputvideopath)
    concat_clip = AudioFileClip(convertfile)
    final_clip = video_clip.set_audio(concat_clip)
    final_clip.write_videofile(outputpath)

args = sys.argv
if 3 <= len(args):
    tagetfile = args[1]
    outputpath = args[2]
    slide_time = int(args[3])
    print(tagetfile)
    print(outputpath)
    with tempfile.TemporaryDirectory() as dname:
        jpg2video(tagetfile,dname,outputpath,slide_time)    
else:
    print('Arguments are too short')

